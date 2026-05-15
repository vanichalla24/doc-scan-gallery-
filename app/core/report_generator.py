"""Report generation: PPTX, Excel, CSV, HTML for TransLingo QA Studio."""
from __future__ import annotations

import csv
import io
import json
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from loguru import logger
except ImportError:
    import logging
    logger = logging.getLogger(__name__)

from app.models.data_models import EngineResult, ScoringEngine, ValidationRun


class ReportGenerator:
    """Generates PPTX, Excel, CSV, and HTML reports from a ValidationRun."""

    def __init__(self, run: ValidationRun):
        self.run = run

    # ── PPTX ────────────────────────────────────────────────────────────────

    def generate_pptx(self, output_path: str) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError as e:
            raise RuntimeError(f"python-pptx or matplotlib not installed: {e}")

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank = prs.slide_layouts[6]

        def add_title_slide():
            slide = prs.slides.add_slide(blank)
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = RGBColor(0x1e, 0x29, 0x3b)

            txb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
            tf = txb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "TransLingo QA Studio"
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

            sub = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
            stf = sub.text_frame
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            sr = sp.add_run()
            sr.text = (
                f"Validation Report  |  "
                f"{self.run.source_language} → {self.run.target_language}  |  "
                f"{self.run.start_time.strftime('%Y-%m-%d')}"
            )
            sr.font.size = Pt(18)
            sr.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)

        def add_summary_slide():
            slide = prs.slides.add_slide(blank)
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = RGBColor(0x1e, 0x29, 0x3b)

            def add_text(x, y, w, h, text, size=16, bold=False, color=RGBColor(0xff, 0xff, 0xff)):
                txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                tf = txb.text_frame
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = text
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = color

            add_text(0.5, 0.3, 12, 0.6, "Executive Summary", 28, True)
            stats = self.run.get_overall_stats()
            add_text(0.5, 1.2, 4, 0.5, f"Total Images: {stats.get('total', 0)}", 16)
            add_text(0.5, 1.8, 4, 0.5, f"Pass Rate: {stats.get('pass_rate', 0):.1f}%", 16)
            add_text(0.5, 2.4, 4, 0.5, f"Avg Score: {stats.get('average_score', 0):.1f}", 16)
            add_text(0.5, 3.0, 4, 0.5, f"Engines Tested: {len(self.run.engines)}", 16)

            y = 1.2
            add_text(5, 1.0, 7, 0.5, "Engine Rankings", 18, True)
            engine_avgs = stats.get("engine_averages", {})
            for eng, avg in sorted(engine_avgs.items(), key=lambda x: -x[1]):
                add_text(5, y, 7, 0.5, f"{eng}: {avg:.1f}", 15)
                y += 0.55

        def add_engine_chart_slide():
            slide = prs.slides.add_slide(blank)
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = RGBColor(0x1e, 0x29, 0x3b)

            stats = self.run.get_overall_stats()
            engine_avgs = stats.get("engine_averages", {})
            if not engine_avgs:
                return

            fig, ax = plt.subplots(figsize=(9, 4), facecolor="#1e293b")
            ax.set_facecolor("#1e293b")
            engines = list(engine_avgs.keys())
            scores = list(engine_avgs.values())
            colors = ["#22c55e" if s >= 85 else "#f59e0b" if s >= 70 else "#ef4444" for s in scores]
            bars = ax.bar(engines, scores, color=colors, edgecolor="#334155")
            ax.set_ylim(0, 100)
            ax.set_ylabel("Score", color="white")
            ax.set_title("Engine Average Scores", color="white", fontsize=14)
            ax.tick_params(colors="white")
            for spine in ax.spines.values():
                spine.set_edgecolor("#334155")
            for bar, score in zip(bars, scores):
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1,
                        f"{score:.1f}", ha="center", color="white", fontsize=11)

            buf = io.BytesIO()
            plt.tight_layout()
            plt.savefig(buf, format="png", dpi=120, facecolor="#1e293b")
            plt.close(fig)
            buf.seek(0)

            txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
            tf = txb.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "Engine Comparison"
            r.font.size = Pt(24)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

            slide.shapes.add_picture(buf, Inches(2), Inches(1.2), Inches(9), Inches(4.5))

        def add_image_slide(result):
            slide = prs.slides.add_slide(blank)
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = RGBColor(0x0f, 0x17, 0x2a)

            def txt(x, y, w, h, t, sz=12, bold=False, color=RGBColor(0xff, 0xff, 0xff)):
                txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                tf = txb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = t
                r.font.size = Pt(sz)
                r.font.bold = bold
                r.font.color.rgb = color

            band_colors = {
                "Excellent": RGBColor(0x22, 0xc5, 0x5e),
                "Good": RGBColor(0x3b, 0x82, 0xf6),
                "Acceptable": RGBColor(0xf5, 0x9e, 0x0b),
                "Needs Review": RGBColor(0xef, 0x44, 0x44),
            }
            band_color = band_colors.get(result.get("score_band", ""), RGBColor(0x94, 0xa3, 0xb8))

            txt(0.3, 0.15, 9, 0.5, f"{result['image_name']} — {result['engine']}", 16, True)
            txt(9.5, 0.15, 3, 0.5,
                f"Score: {result['overall_score']:.1f}  [{result.get('score_band', '')}]",
                14, True, band_color)

            orig_path = result.get("original_path", "")
            trans_path = result.get("translated_path", "")

            if orig_path and Path(orig_path).exists():
                try:
                    slide.shapes.add_picture(orig_path, Inches(0.3), Inches(0.9), Inches(4.5), Inches(3.5))
                    txt(0.3, 4.5, 4.5, 0.4, "Original", 11, color=RGBColor(0x94, 0xa3, 0xb8))
                except Exception:
                    pass

            if trans_path and Path(trans_path).exists():
                try:
                    slide.shapes.add_picture(trans_path, Inches(5.2), Inches(0.9), Inches(4.5), Inches(3.5))
                    txt(5.2, 4.5, 4.5, 0.4, f"Translated ({result['engine']})", 11,
                        color=RGBColor(0x94, 0xa3, 0xb8))
                except Exception:
                    pass

            scores_text = ""
            for k, v in result.get("parameter_scores", {}).items():
                scores_text += f"{v.get('score', 0):.0f}  {k.replace('_', ' ').title()}\n"

            txt(10.0, 0.9, 3.1, 3.5, scores_text or "No scores", 10)

            issues = result.get("issues", [])
            if issues:
                txt(0.3, 5.1, 12.5, 0.4, "Issues:", 10, True, RGBColor(0xf5, 0x9e, 0x0b))
                txt(0.3, 5.5, 12.5, 1.5, "  •  ".join(issues[:5]), 9,
                    color=RGBColor(0xfc, 0xa5, 0xa5))

        add_title_slide()
        add_summary_slide()
        add_engine_chart_slide()

        for engine_name, er in self.run.engine_results.items():
            for img_result in sorted(er.image_results, key=lambda r: r.overall_score):
                d = {
                    "image_name": img_result.image_name,
                    "engine": img_result.engine,
                    "original_path": img_result.original_path,
                    "translated_path": img_result.translated_path,
                    "overall_score": img_result.overall_score,
                    "score_band": img_result.score_band,
                    "parameter_scores": {
                        k: {"score": v.score} for k, v in img_result.parameter_scores.items()
                    },
                    "issues": img_result.issues,
                }
                add_image_slide(d)

        prs.save(output_path)
        logger.info(f"PPTX saved to {output_path}")
        return output_path

    # ── Excel ────────────────────────────────────────────────────────────────

    def generate_excel(self, output_path: str) -> str:
        try:
            import pandas as pd
            from openpyxl.styles import PatternFill, Font, Alignment
            from openpyxl.utils import get_column_letter
        except ImportError as e:
            raise RuntimeError(f"pandas/openpyxl not installed: {e}")

        rows = []
        for engine_name, er in self.run.engine_results.items():
            for r in er.image_results:
                row = {
                    "Engine": engine_name,
                    "Image": r.image_name,
                    "Overall Score": r.overall_score,
                    "Score Band": r.score_band,
                    "Processing Time (s)": round(r.processing_time, 2),
                    "Error": r.error or "",
                }
                for param_name, ps in r.parameter_scores.items():
                    row[param_name.replace("_", " ").title()] = ps.score
                row["Issues"] = " | ".join(r.issues[:5])
                rows.append(row)

        df = pd.DataFrame(rows)

        summary_rows = []
        for engine_name, er in self.run.engine_results.items():
            summary_rows.append({
                "Engine": engine_name,
                "Average Score": round(er.average_score, 2),
                "Pass Rate (%)": round(er.pass_rate, 2),
                "Total Images": er.total_images,
                "Processed": er.processed_images,
            })
        df_summary = pd.DataFrame(summary_rows)

        with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
            df_summary.to_excel(writer, sheet_name="Summary", index=False)
            df.to_excel(writer, sheet_name="Detailed Results", index=False)

            ws = writer.sheets["Detailed Results"]
            for col in ws.columns:
                max_len = max(len(str(cell.value or "")) for cell in col)
                ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 2, 40)

            fill_map = {
                "Excellent": "22c55e",
                "Good": "3b82f6",
                "Acceptable": "f59e0b",
                "Needs Review": "ef4444",
            }
            score_band_col = None
            for idx, cell in enumerate(ws[1], 1):
                if cell.value == "Score Band":
                    score_band_col = idx
                    break

            if score_band_col:
                for row in ws.iter_rows(min_row=2, max_row=ws.max_row):
                    cell = row[score_band_col - 1]
                    band = cell.value or ""
                    hex_color = fill_map.get(band, "6b7280")
                    cell.fill = PatternFill(start_color=hex_color, end_color=hex_color, fill_type="solid")
                    cell.font = Font(color="FFFFFF", bold=True)

        logger.info(f"Excel saved to {output_path}")
        return output_path

    # ── CSV ──────────────────────────────────────────────────────────────────

    def generate_csv(self, output_path: str) -> str:
        rows = []
        headers = None

        for engine_name, er in self.run.engine_results.items():
            for r in er.image_results:
                row = {
                    "engine": engine_name,
                    "image": r.image_name,
                    "overall_score": r.overall_score,
                    "score_band": r.score_band,
                    "processing_time": round(r.processing_time, 2),
                }
                for k, ps in r.parameter_scores.items():
                    row[k] = ps.score
                row["issues"] = " | ".join(r.issues[:3])
                row["error"] = r.error or ""
                if headers is None:
                    headers = list(row.keys())
                rows.append(row)

        with open(output_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.DictWriter(f, fieldnames=headers or [])
            writer.writeheader()
            writer.writerows(rows)

        logger.info(f"CSV saved to {output_path}")
        return output_path

    # ── HTML ─────────────────────────────────────────────────────────────────

    def generate_html(self, output_path: str) -> str:
        stats = self.run.get_overall_stats()
        engine_avgs = stats.get("engine_averages", {})

        def band_class(score):
            if score >= 95:
                return "excellent"
            elif score >= 85:
                return "good"
            elif score >= 70:
                return "acceptable"
            return "needs-review"

        engine_rows = ""
        for eng, avg in sorted(engine_avgs.items(), key=lambda x: -x[1]):
            cls = band_class(avg)
            engine_rows += f'<tr><td>{eng}</td><td class="score {cls}">{avg:.1f}</td></tr>\n'

        detail_rows = ""
        for engine_name, er in self.run.engine_results.items():
            for r in er.image_results:
                cls = band_class(r.overall_score)
                issues_html = "<br>".join(r.issues[:3]) if r.issues else "—"
                detail_rows += f"""
<tr>
  <td>{engine_name}</td>
  <td>{r.image_name}</td>
  <td class="score {cls}">{r.overall_score:.1f}</td>
  <td class="{cls}">{r.score_band}</td>
  <td>{issues_html}</td>
</tr>"""

        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>TransLingo QA Report</title>
<style>
  body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
         background:#0f172a; color:#e2e8f0; margin:0; padding:24px; }}
  h1 {{ color:#f8fafc; font-size:2rem; margin-bottom:4px; }}
  h2 {{ color:#94a3b8; font-size:1.1rem; margin-bottom:32px; }}
  h3 {{ color:#cbd5e1; margin-top:32px; }}
  .cards {{ display:flex; gap:16px; flex-wrap:wrap; margin:24px 0; }}
  .card {{ background:#1e293b; border-radius:12px; padding:20px 28px;
           min-width:160px; }}
  .card .num {{ font-size:2rem; font-weight:700; color:#f8fafc; }}
  .card .lbl {{ color:#64748b; font-size:.85rem; margin-top:4px; }}
  table {{ width:100%; border-collapse:collapse; background:#1e293b;
           border-radius:12px; overflow:hidden; margin-top:16px; }}
  th {{ background:#0f172a; color:#94a3b8; padding:10px 14px;
        text-align:left; font-size:.8rem; text-transform:uppercase; }}
  td {{ padding:9px 14px; border-bottom:1px solid #334155; font-size:.85rem; }}
  tr:last-child td {{ border-bottom:none; }}
  .score {{ font-weight:700; }}
  .excellent {{ color:#22c55e; }}
  .good {{ color:#3b82f6; }}
  .acceptable {{ color:#f59e0b; }}
  .needs-review {{ color:#ef4444; }}
</style>
</head>
<body>
<h1>TransLingo QA Studio</h1>
<h2>{self.run.source_language} → {self.run.target_language} &nbsp;|&nbsp;
    {self.run.start_time.strftime('%Y-%m-%d %H:%M')}</h2>

<div class="cards">
  <div class="card"><div class="num">{stats.get('total',0)}</div><div class="lbl">Total Images</div></div>
  <div class="card"><div class="num">{stats.get('pass_rate',0):.1f}%</div><div class="lbl">Pass Rate</div></div>
  <div class="card"><div class="num">{stats.get('average_score',0):.1f}</div><div class="lbl">Avg Score</div></div>
  <div class="card"><div class="num">{len(self.run.engines)}</div><div class="lbl">Engines</div></div>
</div>

<h3>Engine Rankings</h3>
<table>
  <tr><th>Engine</th><th>Average Score</th></tr>
  {engine_rows}
</table>

<h3>Detailed Results</h3>
<table>
  <tr><th>Engine</th><th>Image</th><th>Score</th><th>Band</th><th>Issues</th></tr>
  {detail_rows}
</table>
</body>
</html>"""

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)

        logger.info(f"HTML report saved to {output_path}")
        return output_path
